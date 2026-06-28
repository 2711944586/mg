from __future__ import annotations

from pathlib import Path
from typing import Iterable

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
PHOTOS = ROOT / "WebDemo" / "public" / "photos"
PUBLIC = ROOT / "WebDemo" / "public"
DOCS = ROOT / "docs"
DELIVERABLES = ROOT / "deliverables"

PPTX_OUT = DELIVERABLES / "守正创新经典品读汇报.pptx"
DOCX_OUT = DELIVERABLES / "汇报讲稿.docx"

RED_DARK = "3A0309"
RED = "8F1014"
RED_BRIGHT = "C82020"
PAPER = "FFF3D6"
PAPER_DEEP = "E8C983"
INK = "34100A"
GOLD = "E8B85A"
MUTED = "6D3A21"
WHITE = "FFF8EA"


def rgb(hex_value: str) -> PptRGBColor:
    return PptRGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


slides = [
    {
        "kicker": "经典品读汇报",
        "title": "守正创新何以开辟新境界",
        "subtitle": "从“两个结合”到新时代青年的理论自觉与实践担当",
        "time": "20秒",
        "image": PHOTOS / "national-library-reading-room.jpg",
        "layout": "cover",
    },
    {
        "kicker": "开卷",
        "title": "带着一个问题进入经典",
        "claim": "为什么马克思主义必须中国化时代化？青年怎样把经典阅读转化为理解时代的能力？",
        "bullets": ["时代之问：回应新的实践课题", "理论之答：守正创新与两个结合", "青年之责：形成方向感、判断力、行动力"],
        "image": PHOTOS / "national-library-reading-room.jpg",
        "time": "45秒",
        "layout": "image_left",
    },
    {
        "kicker": "主线",
        "title": "理论在回答中国问题中生长",
        "claim": "马克思主义中国化时代化贯穿革命、建设、改革和新时代。",
        "timeline": [
            ("原理", "马克思主义基本原理"),
            ("革命", "毛泽东思想"),
            ("改革", "中国特色社会主义理论体系"),
            ("新时代", "习近平新时代中国特色社会主义思想"),
            ("继续", "推进中国化时代化"),
        ],
        "time": "60秒",
        "layout": "timeline",
    },
    {
        "kicker": "方法",
        "title": "守正创新不是两选一",
        "claim": "守正保证根本方向，创新回应新的问题，二者统一于实践。",
        "duo": [("守正", "坚持马克思主义基本立场、观点、方法"), ("创新", "回应中国之问、世界之问、人民之问、时代之问")],
        "quote": "守正不是守旧，创新不是离根。",
        "image": PHOTOS / "ten-bamboo-page.jpg",
        "time": "60秒",
        "layout": "duo",
    },
    {
        "kicker": "路径一",
        "title": "同中国具体实际相结合",
        "claim": "第一个结合回答理论如何不悬空。",
        "bullets": ["从中国国情出发", "回答中国自己的问题", "不照搬别国模式", "对应实事求是、独立自主"],
        "image": PHOTOS / "wuyuan-village.jpg",
        "time": "65秒",
        "layout": "image_right",
    },
    {
        "kicker": "路径二",
        "title": "同中华优秀传统文化相结合",
        "claim": "第二个结合回答理论如何有根脉。",
        "bullets": ["在中华文明土壤中扎根", "让优秀传统文化提供滋养", "推动创造性转化、创新性发展"],
        "tags": ["民为邦本", "天下为公", "和而不同", "自强不息"],
        "image": PHOTOS / "forbidden-city.jpg",
        "time": "65秒",
        "layout": "image_left",
    },
    {
        "kicker": "现实",
        "title": "问题导向让理论进入现场",
        "claim": "理论生命力体现在解释现实、回应现实、指导现实。",
        "cases": [
            ("科技自立自强", "时代化与新质生产力", PHOTOS / "fast-telescope.jpg"),
            ("乡村振兴", "中国实际与共同富裕", PHOTOS / "jiangxi-village.jpg"),
            ("绿色发展", "生态文明与现代化", PHOTOS / "saihanba-forest.jpg"),
            ("交通强国", "独立自主与道路选择", PHOTOS / "high-speed-rail.jpg"),
        ],
        "time": "70秒",
        "layout": "cases",
    },
    {
        "kicker": "共读",
        "title": "经典阅读不是背概念",
        "claim": "共读让文本从个人理解走向共同表达。",
        "steps": [("读原文", "圈画关键词，梳理文献结构"), ("辨问题", "把文本放回毛概课程"), ("连现实", "用现实案例检验解释力"), ("定表达", "形成PPT、讲稿、网页和短片")],
        "time": "50秒",
        "layout": "steps",
    },
    {
        "kicker": "展厅",
        "title": "WebDemo把全篇逻辑翻成一本书",
        "claim": "自动翻书演示把12页PPT压缩成九章总结，适合投屏、停留和复盘。",
        "bullets": ["红色动态背景营造课堂仪式感", "真实图像连接阅读、根脉和现实中国", "线下讨论板块承接小组共读现场"],
        "image": PUBLIC / "reading-hero.png",
        "time": "60秒",
        "layout": "image_right",
    },
    {
        "kicker": "短片",
        "title": "让理论之光照亮青春之路",
        "claim": "从书页到时代场景，从关键词到现实案例，从经典阅读到青年担当。",
        "steps": [("书页", "翻开经典文本"), ("关键词", "守正、创新、两个结合"), ("现实", "城市、乡村、科技、生态"), ("青年", "阅读、讨论、实践")],
        "time": "40秒",
        "layout": "steps",
    },
    {
        "kicker": "青年",
        "title": "从理论自觉到实践担当",
        "claim": "把个人成长放进国家需要和时代趋势中思考。",
        "bullets": ["学理论：不是停留在背诵，而是掌握方法", "看时代：不是旁观变化，而是理解方向", "做青年：不是只谈理想，而是提升本领、服务社会"],
        "image": PHOTOS / "national-library-hall.jpg",
        "time": "55秒",
        "layout": "image_left",
    },
    {
        "kicker": "合卷",
        "title": "在经典中坚定方向，在实践中增长本领",
        "claim": "守正创新、两个结合、青年担当，构成这次经典品读的完整闭环。",
        "summary": [("守正创新", "开辟新境界的方法"), ("两个结合", "理论扎根中国、回应时代的路径"), ("青年担当", "经典阅读最终落到的行动")],
        "time": "40秒",
        "layout": "summary",
    },
]


def set_text_style(paragraph, size=18, color=INK, bold=False, font="Microsoft YaHei", align=None):
    if align is not None:
        paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = font
        run.font.size = PptPt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)


def add_textbox(slide, text, x, y, w, h, size=18, color=INK, bold=False, font="Microsoft YaHei", align=None, fill=None):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = PptInches(0.06)
    tf.margin_right = PptInches(0.06)
    tf.margin_top = PptInches(0.04)
    tf.margin_bottom = PptInches(0.04)
    p = tf.paragraphs[0]
    p.text = text
    set_text_style(p, size=size, color=color, bold=bold, font=font, align=align)
    if fill:
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(fill)
        box.line.color.rgb = rgb(fill)
    return box


def add_rect(slide, x, y, w, h, fill, transparency=0, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.fill.transparency = transparency
    if line:
        shp.line.color.rgb = rgb(line)
    else:
        shp.line.fill.background()
    return shp


def add_picture_cover(slide, path: Path, x, y, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    target_ratio = w / h
    source_ratio = iw / ih
    if source_ratio > target_ratio:
        pic_h = h
        pic_w = h * source_ratio
        px = x - (pic_w - w) / 2
        py = y
    else:
        pic_w = w
        pic_h = w / source_ratio
        px = x
        py = y - (pic_h - h) / 2
    pic = slide.shapes.add_picture(str(path), PptInches(px), PptInches(py), PptInches(pic_w), PptInches(pic_h))
    return pic


def send_to_back(slide, shape):
    sp_tree = slide.shapes._spTree
    sp_tree.remove(shape._element)
    sp_tree.insert(2, shape._element)


def background(slide):
    bg = add_rect(slide, 0, 0, 13.333, 7.5, RED_DARK)
    send_to_back(slide, bg)
    add_rect(slide, 0, 0, 13.333, 7.5, "180105", 0)
    add_rect(slide, 0.28, 0.25, 12.78, 7.0, PAPER, 0, line="DDBE76")
    add_rect(slide, 0.52, 0.48, 12.3, 6.55, "FFF8EA", 0, line="E6C77B")


def title_block(slide, data):
    add_textbox(slide, data["kicker"], 0.74, 0.62, 2.0, 0.28, size=10, color=RED, bold=True)
    add_textbox(slide, data["title"], 0.74, 0.9, 8.7, 0.82, size=28, color=INK, bold=True, font="SimSun")
    add_textbox(slide, data["time"], 11.4, 0.72, 0.86, 0.34, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, fill=RED)


def add_claim(slide, claim, x=0.82, y=1.78, w=5.3, h=0.75):
    add_textbox(slide, claim, x, y, w, h, size=15, color=MUTED, bold=True, font="Microsoft YaHei")


def slide_cover(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_cover(slide, data["image"], 0, 0, 13.333, 7.5)
    add_rect(slide, 0, 0, 13.333, 7.5, "160105", 24)
    add_rect(slide, 0.7, 0.62, 11.9, 6.25, "160105", 26, line="E3BC62")
    add_textbox(slide, data["kicker"], 1.05, 0.98, 2.8, 0.34, size=12, color=GOLD, bold=True)
    add_textbox(slide, data["title"], 1.02, 1.65, 8.6, 1.22, size=39, color=WHITE, bold=True, font="SimSun")
    add_textbox(slide, data["subtitle"], 1.08, 3.02, 7.4, 0.72, size=18, color="F4DCA5", bold=True)
    add_textbox(slide, "《开辟马克思主义中国化时代化新境界》经典品读汇报", 1.08, 4.22, 6.8, 0.42, size=14, color=WHITE)
    add_textbox(slide, "10分30秒课堂展示", 1.08, 5.82, 2.2, 0.36, size=12, color=RED_DARK, bold=True, align=PP_ALIGN.CENTER, fill=GOLD)


def slide_image_left(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    add_picture_cover(slide, data["image"], 0.74, 1.72, 5.2, 4.9)
    add_rect(slide, 0.74, 1.72, 5.2, 4.9, RED_DARK, 72, line="DDBE76")
    title_block(slide, data)
    add_claim(slide, data["claim"], 6.35, 1.82, 5.55, 0.88)
    y = 2.92
    for idx, bullet in enumerate(data.get("bullets", []), 1):
        add_textbox(slide, f"{idx:02d}", 6.42, y + 0.05, 0.5, 0.34, size=11, color=GOLD, bold=True, align=PP_ALIGN.CENTER, fill=RED)
        add_textbox(slide, bullet, 7.02, y, 4.8, 0.44, size=15, color=INK, bold=True)
        y += 0.74
    if "tags" in data:
        tx = 6.42
        for tag in data["tags"]:
            add_textbox(slide, tag, tx, 5.78, 1.12, 0.34, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER, fill=RED)
            tx += 1.24


def slide_image_right(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    add_picture_cover(slide, data["image"], 7.36, 1.68, 4.95, 4.92)
    add_rect(slide, 7.36, 1.68, 4.95, 4.92, RED_DARK, 72, line="DDBE76")
    title_block(slide, data)
    add_claim(slide, data["claim"], 0.82, 1.86, 5.7, 0.82)
    y = 2.98
    for idx, bullet in enumerate(data.get("bullets", []), 1):
        add_textbox(slide, f"{idx:02d}", 0.9, y + 0.04, 0.5, 0.34, size=11, color=GOLD, bold=True, align=PP_ALIGN.CENTER, fill=RED)
        add_textbox(slide, bullet, 1.52, y, 4.8, 0.44, size=15, color=INK, bold=True)
        y += 0.72


def slide_timeline(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    title_block(slide, data)
    add_claim(slide, data["claim"], 0.82, 1.82, 8.1, 0.6)
    base_y = 3.42
    add_rect(slide, 1.1, base_y + 0.39, 10.9, 0.05, RED)
    x = 1.0
    for idx, (year, label) in enumerate(data["timeline"], 1):
        add_rect(slide, x, base_y, 1.62, 0.82, RED if idx % 2 else "6D2A1B", 0, radius=True)
        add_textbox(slide, year, x + 0.12, base_y + 0.1, 1.38, 0.26, size=11, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, x + 0.12, base_y + 0.38, 1.38, 0.34, size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += 2.25
    add_textbox(slide, "理论总是在实践提出的新问题中获得新的表达。", 1.28, 5.35, 10.4, 0.52, size=20, color=INK, bold=True, font="SimSun", align=PP_ALIGN.CENTER)


def slide_duo(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    add_picture_cover(slide, data["image"], 8.3, 1.72, 3.78, 4.84)
    add_rect(slide, 8.3, 1.72, 3.78, 4.84, RED_DARK, 66, line="DDBE76")
    title_block(slide, data)
    add_claim(slide, data["claim"], 0.82, 1.8, 6.3, 0.7)
    x = 0.9
    for name, desc in data["duo"]:
        add_rect(slide, x, 2.92, 3.18, 1.66, "FFF9EA", 0, line="DDBE76", radius=True)
        add_textbox(slide, name, x + 0.22, 3.14, 2.72, 0.42, size=22, color=RED, bold=True, font="SimSun", align=PP_ALIGN.CENTER)
        add_textbox(slide, desc, x + 0.34, 3.74, 2.52, 0.48, size=13, color=INK, bold=True, align=PP_ALIGN.CENTER)
        x += 3.58
    add_textbox(slide, data["quote"], 1.2, 5.35, 6.4, 0.54, size=22, color=WHITE, bold=True, font="SimSun", align=PP_ALIGN.CENTER, fill=RED)


def slide_cases(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    title_block(slide, data)
    add_claim(slide, data["claim"], 0.82, 1.78, 8.6, 0.58)
    positions = [(0.86, 2.55), (6.62, 2.55), (0.86, 4.78), (6.62, 4.78)]
    for (title, note, img), (x, y) in zip(data["cases"], positions):
        add_picture_cover(slide, img, x, y, 5.25, 1.72)
        add_rect(slide, x, y, 5.25, 1.72, RED_DARK, 62, line="DDBE76")
        add_textbox(slide, title, x + 0.24, y + 0.92, 2.45, 0.34, size=15, color=WHITE, bold=True, font="SimSun")
        add_textbox(slide, note, x + 0.24, y + 1.26, 3.2, 0.28, size=10.5, color="F6D995", bold=True)


def slide_steps(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    title_block(slide, data)
    add_claim(slide, data["claim"], 0.82, 1.78, 9.2, 0.6)
    x = 0.94
    for idx, (name, desc) in enumerate(data["steps"], 1):
        add_rect(slide, x, 3.0, 2.65, 1.86, "FFF9EA", 0, line="DDBE76", radius=True)
        add_textbox(slide, f"{idx:02d}", x + 0.2, 3.2, 0.55, 0.34, size=12, color=GOLD, bold=True, align=PP_ALIGN.CENTER, fill=RED)
        add_textbox(slide, name, x + 0.86, 3.16, 1.4, 0.36, size=18, color=RED, bold=True, font="SimSun")
        add_textbox(slide, desc, x + 0.32, 3.82, 2.0, 0.5, size=12.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
        x += 3.0


def slide_summary(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, RED_DARK)
    add_rect(slide, 0.48, 0.44, 12.36, 6.62, "220106", 10, line=GOLD)
    add_textbox(slide, data["kicker"], 0.9, 0.82, 1.4, 0.28, size=11, color=GOLD, bold=True)
    add_textbox(slide, data["title"], 0.9, 1.28, 8.8, 1.0, size=32, color=WHITE, bold=True, font="SimSun")
    add_textbox(slide, data["claim"], 0.94, 2.3, 8.8, 0.58, size=16, color="F6D995", bold=True)
    x = 0.98
    for idx, (title, note) in enumerate(data["summary"], 1):
        add_rect(slide, x, 3.45, 3.62, 1.34, "FFF4D8", 0, line=GOLD, radius=True)
        add_textbox(slide, f"{idx:02d}", x + 0.2, 3.72, 0.5, 0.32, size=12, color=GOLD, bold=True, align=PP_ALIGN.CENTER, fill=RED)
        add_textbox(slide, title, x + 0.86, 3.62, 2.0, 0.38, size=18, color=RED, bold=True, font="SimSun")
        add_textbox(slide, note, x + 0.86, 4.1, 2.1, 0.34, size=11.5, color=INK, bold=True)
        x += 3.94
    add_textbox(slide, "读经典，不是停留在文本表面，而是在理论中看清时代，在实践中坚定方向。", 1.1, 5.72, 11.2, 0.46, size=18, color=WHITE, bold=True, font="SimSun", align=PP_ALIGN.CENTER)


def build_pptx():
    DELIVERABLES.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    for data in slides:
        layout = data["layout"]
        if layout == "cover":
            slide_cover(prs, data)
        elif layout == "image_left":
            slide_image_left(prs, data)
        elif layout == "image_right":
            slide_image_right(prs, data)
        elif layout == "timeline":
            slide_timeline(prs, data)
        elif layout == "duo":
            slide_duo(prs, data)
        elif layout == "cases":
            slide_cases(prs, data)
        elif layout == "steps":
            slide_steps(prs, data)
        elif layout == "summary":
            slide_summary(prs, data)
    prs.save(PPTX_OUT)


def set_run_font(run, name="Microsoft YaHei", size=None, color=None, bold=None):
    run.font.name = name
    run._element.rPr.rFonts.set(qn("w:ascii"), name)
    run._element.rPr.rFonts.set(qn("w:hAnsi"), name)
    run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    if size is not None:
        run.font.size = Pt(size)
    if color is not None:
        run.font.color.rgb = RGBColor.from_string(color)
    if bold is not None:
        run.bold = bold


def paragraph_border_bottom(paragraph, color="E8B85A", size="8"):
    p_pr = paragraph._p.get_or_add_pPr()
    p_bdr = p_pr.find(qn("w:pBdr"))
    if p_bdr is None:
        p_bdr = OxmlElement("w:pBdr")
        p_pr.append(p_bdr)
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), size)
    bottom.set(qn("w:space"), "6")
    bottom.set(qn("w:color"), color)
    p_bdr.append(bottom)


def iter_markdown_sections(text: str):
    current = None
    buf: list[str] = []
    for line in text.splitlines():
        if line.startswith("## 第") or line.startswith("## 彩排"):
            if current:
                yield current, "\n".join(buf).strip()
            current = line.replace("## ", "").strip()
            buf = []
        elif current:
            if not line.startswith("---"):
                buf.append(line)
    if current:
        yield current, "\n".join(buf).strip()


def build_docx():
    DELIVERABLES.mkdir(exist_ok=True)
    doc = Document()
    section = doc.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(0.82)
    section.bottom_margin = Inches(0.78)
    section.left_margin = Inches(0.86)
    section.right_margin = Inches(0.86)

    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = "Microsoft YaHei"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    normal.font.size = Pt(10.8)
    normal.paragraph_format.line_spacing = 1.18
    normal.paragraph_format.space_after = Pt(6)

    title = doc.add_paragraph()
    title.paragraph_format.space_after = Pt(2)
    run = title.add_run("守正创新何以开辟新境界")
    set_run_font(run, "SimSun", 24, RED_DARK, True)
    paragraph_border_bottom(title)

    subtitle = doc.add_paragraph()
    subtitle.paragraph_format.space_after = Pt(12)
    run = subtitle.add_run("《开辟马克思主义中国化时代化新境界》经典品读汇报讲稿 | 10分30秒")
    set_run_font(run, "Microsoft YaHei", 11, MUTED, True)

    md = (DOCS / "汇报讲稿.md").read_text(encoding="utf-8")
    for heading, body in iter_markdown_sections(md):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(10)
        p.paragraph_format.space_after = Pt(4)
        run = p.add_run(heading)
        set_run_font(run, "SimSun", 15, RED, True)
        for para in body.split("\n\n"):
            para = para.strip()
            if not para or para.startswith("汇报主题") or para.startswith("预计时长") or para.startswith("建议语速"):
                continue
            if para.startswith("- "):
                for item in para.splitlines():
                    q = doc.add_paragraph(style=None)
                    q.paragraph_format.left_indent = Inches(0.22)
                    q.paragraph_format.first_line_indent = Inches(-0.1)
                    r = q.add_run("• " + item[2:])
                    set_run_font(r, "Microsoft YaHei", 10.5, INK)
                continue
            p2 = doc.add_paragraph()
            p2.paragraph_format.first_line_indent = Inches(0.28)
            p2.paragraph_format.line_spacing = 1.22
            run = p2.add_run(para.replace("\n", ""))
            set_run_font(run, "Microsoft YaHei", 10.8, INK)

    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = footer.add_run("守正创新经典品读汇报")
    set_run_font(run, "Microsoft YaHei", 8.5, "8A6B3A")
    doc.save(DOCX_OUT)


if __name__ == "__main__":
    build_pptx()
    build_docx()
    print(PPTX_OUT)
    print(DOCX_OUT)

